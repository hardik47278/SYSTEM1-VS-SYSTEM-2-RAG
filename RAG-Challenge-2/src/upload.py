# ===============================
# upload.py  (UPLOAD FILE: PDF/DOCX/XLSX/PPTX/CSV)
# LangGraph token streaming enabled using ONLY LangGraph (app.stream)
# SYNC-ONLY VERSION (no async anywhere)
# Fix: MemorySaver checkpointer requires configurable.thread_id
# ===============================

import os
import re
import logging
import hashlib
from uuid import uuid4
from datetime import datetime, timedelta
from pathlib import Path
from typing import List, Tuple
from typing_extensions import TypedDict
from io import BytesIO
import tempfile
from typing import List,Dict,Optional

from dotenv import load_dotenv

from langgraph.graph import StateGraph, END
from langgraph.checkpoint.memory import MemorySaver

from langchain_core.documents import Document as LCDocument
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_core.prompts import PromptTemplate
from langchain_groq import ChatGroq

from unstructured.partition.pdf import partition_pdf
from unstructured.chunking.title import chunk_by_title

from openai import OpenAI
from ragas.llms import llm_factory
from ragas import SingleTurnSample
from ragas.metrics import LLMContextPrecisionWithoutReference

import docx  # python-docx

# ✅ Qdrant wrapper
from vector_db import upsert_embeddings, query_embeddings

# PPTX
from pptx import Presentation
from PIL import Image

# CSVLoader
from langchain_community.document_loaders.csv_loader import CSVLoader

# LlamaIndex (Excel only)
from llama_index.core import VectorStoreIndex, SimpleDirectoryReader, Settings
from llama_index.embeddings.huggingface import HuggingFaceEmbedding
from llama_index.readers.docling import DoclingReader
from llama_index.core.node_parser import MarkdownNodeParser
from llama_index.llms.langchain import LangChainLLM

import pandas as pd

# Audio
import sounddevice as sd
from scipy.io.wavfile import write
import whisper
from elevenlabs import ElevenLabs

from rank_bm25 import BM25Okapi
from richetadataextraction import extract_rich_metadata

_BM25_STORE = {}

# ✅ Absolute base dir — always next to this file, never relative
BASE_DIR = Path(__file__).resolve().parent / "file_indexes"
PPTX_IMAGES_DIR = BASE_DIR / "pptx_images"
BASE_DIR.mkdir(parents=True, exist_ok=True)
PPTX_IMAGES_DIR.mkdir(parents=True, exist_ok=True)

# -------------------------------
# LOGGING
# -------------------------------
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# -------------------------------
# ENV
# -------------------------------
load_dotenv()
GROQ_API_KEY = os.getenv("GROQ_API_KEY")
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
ELEVENLABS_API_KEY = os.getenv("ELEVENLABS_API_KEY")

# -------------------------------
# MODELS
# -------------------------------
EMBEDDINGS = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
LLM = ChatGroq(
    api_key=GROQ_API_KEY,
    model_name="llama-3.1-8b-instant",
    temperature=0,
    streaming=True,
)

# -------------------------------
# CACHE (answers)
# -------------------------------
CACHE = {}
CACHE_TTL = timedelta(hours=1)

# -------------------------------
# CACHE (LlamaIndex query engines) - Excel
# -------------------------------
_LLAMA_CACHE = {}

# ===============================
# WHISPER CACHE
# ===============================
_WHISPER_MODEL = None


def get_whisper_model():
    global _WHISPER_MODEL
    if _WHISPER_MODEL is None:
        logger.info("Loading Whisper model...")
        _WHISPER_MODEL = whisper.load_model("base")
    return _WHISPER_MODEL


# ===============================
# VOICE
# ===============================
def record_user_input(duration=5, samplerate=16000) -> BytesIO:
    recording = sd.rec(int(duration * samplerate), samplerate=samplerate, channels=1)
    sd.wait()
    buffer = BytesIO()
    write(buffer, samplerate, recording)
    buffer.seek(0)
    return buffer


def extract_text(audio_buffer: BytesIO) -> str:
    audio_buffer.seek(0)
    with tempfile.NamedTemporaryFile(delete=False, suffix=".wav") as tmp:
        tmp.write(audio_buffer.read())
        tmp_path = tmp.name
    model = get_whisper_model()
    result = model.transcribe(tmp_path)
    os.remove(tmp_path)
    return result["text"]


def generate_audio(text: str, lang="en", voice_id="hpp4J3VqNfWAUOO0d1Us") -> BytesIO:
    client = ElevenLabs(api_key=ELEVENLABS_API_KEY)
    model_id = "eleven_multilingual_v2"
    audio_stream = client.text_to_speech.convert(
        voice_id=voice_id,
        model_id=model_id,
        text=text,
    )
    buffer = BytesIO()
    for chunk in audio_stream:
        if chunk:
            buffer.write(chunk)
    buffer.seek(0)
    return buffer


def make_cache_key(question: str, index_id: str):
    raw = f"{question}|upload_file|{index_id}"
    return hashlib.sha256(raw.encode()).hexdigest()


def cache_get(key: str):
    item = CACHE.get(key)
    if not item:
        return None
    if datetime.now() - item["created_at"] > CACHE_TTL:
        del CACHE[key]
        return None
    return item["answer"], item.get("sources", []), item.get("context_precision", 0.0)


def cache_set(key: str, value: str, sources=None, context_precision: float = 0.0):
    CACHE[key] = {
        "answer": value,
        "sources": sources or [],
        "context_precision": float(context_precision),
        "created_at": datetime.now(),
    }


# -------------------------------
# TEXT CLEANING
# -------------------------------
def normalize_text(s: str) -> str:
    if not s:
        return ""
    s = s.replace("\u00a0", " ")
    s = s.replace("\f", "\n")
    s = re.sub(r"[ \t]+", " ", s)
    s = re.sub(r"\n{3,}", "\n\n", s)
    return s.strip()


def _upsert_lcdocs_to_qdrant(docs: List[LCDocument], file_id: str):
    if not docs:
        return
    texts = [d.page_content for d in docs]
    vectors = EMBEDDINGS.embed_documents(texts)
    payload_vectors = []
    for i, (d, vec) in enumerate(zip(docs, vectors)):
        md = dict(d.metadata or {})
        md["file_id"] = file_id
        md["text"] = d.page_content
        payload_vectors.append(
            {
                
                "values": vec,
                "metadata": md,
            }
        )
    upsert_embeddings(payload_vectors)


# -------------------------------
# CSV INGEST (Qdrant)
def build_documents_from_csvloader(csv_path: str, file_id: str) -> List[LCDocument]:
    start = datetime.now()
    loader = CSVLoader(file_path=csv_path)
    docs = loader.load()

    out: List[LCDocument] = []

    # 1) try CSVLoader docs
    for i, d in enumerate(docs):
        text = normalize_text((d.page_content or "").strip())
        if not text:
            continue

        # allow short rows
        if len(text.split()) < 1:
            continue

        md = dict(d.metadata or {})
        out.append(
            LCDocument(
                page_content=text,
                metadata={
                    "paper_id": file_id,
                    "chunk_id": i,
                    "source_type": "csv",
                    "section_title": f"Row {md.get('row', i)}",
                    **md,
                },
            )
        )

    # 2) fallback: if CSVLoader gave nothing, use pandas
    if not out:
        try:
            df = pd.read_csv(csv_path)
            max_rows = min(len(df), 5000)

            for i in range(max_rows):
                row = df.iloc[i]
                row_text = normalize_text(" | ".join([f"{k}: {row.get(k)}" for k in df.columns]))
                if not row_text:
                    continue

                out.append(
                    LCDocument(
                        page_content=row_text,
                        metadata={
                            "paper_id": file_id,
                            "chunk_id": i,
                            "source_type": "csv",
                            "section_title": f"Row {i}",
                            "row": i,
                        },
                    )
                )
        except Exception as e:
            logger.warning("CSV fallback pandas read failed: %s", e)

    logger.info("⏱ build_documents_from_csvloader: %.3fs", (datetime.now() - start).total_seconds())

    # ✅ DO NOT raise for CSV
    return out


def create_csv_index(csv_path: str, csv_name: str, base_dir: str = None) -> str:
    start = datetime.now()
    docs = build_documents_from_csvloader(csv_path, file_id=csv_name)
    file_id = csv_name

    if docs:
        _upsert_lcdocs_to_qdrant(docs, file_id=file_id)
        logger.info("✅ CSV indexed to Qdrant. chunks=%d", len(docs))
    else:
        logger.warning("⚠️ CSV produced zero docs; skipping Qdrant indexing.")

    logger.info("⏱ create_csv_index TOTAL: %.3fs", (datetime.now() - start).total_seconds())
    return file_id


# -------------------------------
# PDF INGEST (Qdrant)
# -------------------------------
def partition_document_pdf(file_path: str):
    start = datetime.now()
    elements = partition_pdf(filename=file_path, strategy="fast", infer_table_structure=True)
    logger.info("⏱ partition_document_pdf: %.3fs", (datetime.now() - start).total_seconds())
    return elements


def create_chunks_pdf(elements):
    start = datetime.now()
    chunks = chunk_by_title(
        elements,
        max_characters=3000,
        new_after_n_chars=2400,
        combine_text_under_n_chars=500,
    )
    logger.info("⏱ create_chunks_pdf: %.3fs", (datetime.now() - start).total_seconds())
    return chunks


def build_documents_from_pdf_chunks(chunks, paper_id: str) -> List[LCDocument]:
    start = datetime.now()
    docs: List[LCDocument] = []
    for i, chunk in enumerate(chunks):
        text = getattr(chunk, "text", None)
        if not text:
            continue
        text = normalize_text(text)
        if len(text.split()) < 40:
            continue
        docs.append(
            LCDocument(
                page_content=text,
                metadata={"paper_id": paper_id, "chunk_id": i, "source_type": "pdf"},
            )
        )
    logger.info("⏱ build_documents_from_pdf_chunks: %.3fs", (datetime.now() - start).total_seconds())
    if not docs:
        raise RuntimeError("No valid chunks extracted from PDF")
    return docs


# -------------------------------
# DOCX INGEST (Qdrant)
# -------------------------------
def read_docx_sections(path: Path) -> List[Tuple[str, str]]:
    d = docx.Document(str(path))
    sections: List[Tuple[str, str]] = []
    current_title = "Document"
    current_lines: List[str] = []
    for p in d.paragraphs:
        style = (p.style.name or "").lower()
        txt = normalize_text(p.text)
        if not txt:
            continue
        if "heading" in style:
            if current_lines:
                sections.append((current_title, "\n".join(current_lines).strip()))
                current_lines = []
            current_title = txt
        else:
            current_lines.append(txt)
    if current_lines:
        sections.append((current_title, "\n".join(current_lines).strip()))
    sections = [(t, c) for (t, c) in sections if c and len(c.split()) >= 40]
    if not sections:
        raise RuntimeError("No valid sections extracted from DOCX (too little text?)")
    return sections


def build_documents_from_docx_sections(sections: List[Tuple[str, str]], file_id: str) -> List[LCDocument]:
    docs: List[LCDocument] = []
    for i, (title, content) in enumerate(sections):
        docs.append(
            LCDocument(
                page_content=f"{title}\n\n{content}",
                metadata={"paper_id": file_id, "chunk_id": i, "source_type": "docx", "section_title": title},
            )
        )
    return docs


# -------------------------------
# PPTX INGEST (Qdrant)
# -------------------------------
def build_documents_from_pptx(
    pptx_path: str,
    file_id: str,
    images_output_dir: str = None,   # ✅ now optional — defaults to absolute path
    save_images: bool = True,
) -> List[LCDocument]:
    start = datetime.now()

    # ✅ Always use absolute path next to upload.py
    if images_output_dir is None:
        images_output_dir = str(PPTX_IMAGES_DIR)

    prs = Presentation(str(pptx_path))

    if save_images:
        Path(images_output_dir).mkdir(parents=True, exist_ok=True)

    docs: List[LCDocument] = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        full_text: List[str] = []
        extracted_images: List[str] = []

        full_text.append(f"### Slide {slide_num} ###")

        for shape_num, shape in enumerate(slide.shapes, start=1):
            if hasattr(shape, "text") and shape.text and shape.text.strip():
                full_text.append(shape.text.strip())

            if getattr(shape, "has_table", False):
                full_text.append(f"--- Table in Slide {slide_num}, Shape {shape_num} ---")
                table = shape.table
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    full_text.append(" | ".join(row_data))

            if save_images and hasattr(shape, "image") and shape.image is not None:
                try:
                    image_bytes = shape.image.blob
                    img_pil = Image.open(BytesIO(image_bytes))
                    img_filename = f"{Path(file_id).stem}_slide{slide_num}_img{shape_num}.png"
                    img_path = str(Path(images_output_dir) / img_filename)
                    img_pil.save(img_path, format="PNG")
                    extracted_images.append(img_path)
                except Exception as e:
                    logger.warning("Error extracting image from shape: %s", e)

        slide_text = "\n".join(full_text).strip()
        if len(slide_text.split()) < 10:
            continue

        docs.append(
            LCDocument(
                page_content=slide_text,
                metadata={
                    "paper_id": file_id,
                    "chunk_id": slide_num,
                    "source_type": "pptx",
                    "section_title": f"Slide {slide_num}",
                    "slide_num": slide_num,
                    "image_count": len(extracted_images),
                    "images": extracted_images,
                },
            )
        )

    logger.info("⏱ build_documents_from_pptx: %.3fs", (datetime.now() - start).total_seconds())
    if not docs:
        raise RuntimeError("No valid chunks extracted from PPTX (too little text?)")
    return docs


# -------------------------------
# INDEX BUILDERS
# -------------------------------
def create_pdf_index(pdf_path: str, pdf_name: str, base_dir: str = None) -> str:
    start = datetime.now()
    elements = partition_document_pdf(pdf_path)
    chunks = create_chunks_pdf(elements)
    docs = build_documents_from_pdf_chunks(chunks, paper_id=pdf_name)
    try:
        combined_text = "\n\n".join([d.page_content for d in docs[:5]])
        rich_meta = extract_rich_metadata(combined_text)
        for d in docs:
            d.metadata.update({
                "title": rich_meta.get("title"),
                "author": rich_meta.get("author"),
                "summary": rich_meta.get("summary"),
                "keywords": rich_meta.get("keywords"),
                "year": rich_meta.get("year"),
            })
        logger.info("rich metadata extracted successfully")
    except Exception as e:
        logger.warning("metadata extraction failed: %s", e)

    file_id = pdf_name
    _upsert_lcdocs_to_qdrant(docs, file_id=file_id)
    logger.info("⏱ create_pdf_index TOTAL: %.3fs", (datetime.now() - start).total_seconds())
    return file_id


def create_docx_index(docx_path: str, docx_name: str, base_dir: str = None) -> str:
    start = datetime.now()
    sections = read_docx_sections(Path(docx_path))
    docs = build_documents_from_docx_sections(sections, file_id=docx_name)
    file_id = docx_name
    _upsert_lcdocs_to_qdrant(docs, file_id=file_id)
    logger.info("⏱ create_docx_index TOTAL: %.3fs", (datetime.now() - start).total_seconds())
    return file_id


def create_pptx_index(pptx_path: str, pptx_name: str, base_dir: str = None) -> str:
    start = datetime.now()

    # ✅ Always use absolute path — no more permission errors
    docs = build_documents_from_pptx(
        pptx_path=pptx_path,
        file_id=pptx_name,
        images_output_dir=str(PPTX_IMAGES_DIR),
        save_images=True,
    )

    file_id = pptx_name
    _upsert_lcdocs_to_qdrant(docs, file_id=file_id)
    logger.info("⏱ create_pptx_index TOTAL: %.3fs", (datetime.now() - start).total_seconds())
    return file_id


def create_excel_llama_index(file_path: str, filename: str) -> str:
    start = datetime.now()
    reader = DoclingReader()
    loader = SimpleDirectoryReader(
        input_files=[file_path],
        file_extractor={".xlsx": reader, ".xls": reader},
    )
    docs = loader.load_data()
    Settings.embed_model = HuggingFaceEmbedding(model_name="sentence-transformers/all-MiniLM-L6-v2")
    Settings.llm = LangChainLLM(LLM)
    node_parser = MarkdownNodeParser()
    index = VectorStoreIndex.from_documents(
        documents=docs,
        transformations=[node_parser],
        show_progress=True,
    )
    query_engine = index.as_query_engine(streaming=False)
    _LLAMA_CACHE[filename] = query_engine
    logger.info("⏱ create_excel_llama_index TOTAL: %.3fs", (datetime.now() - start).total_seconds())
    return filename


def create_file_index(file_path: str, filename: str, base_dir: str = None) -> str:
    suffix = Path(filename).suffix.lower()

    if suffix == ".pdf":
        return create_pdf_index(file_path, filename)
    if suffix == ".docx":
        return create_docx_index(file_path, filename)
    if suffix == ".pptx":
        return create_pptx_index(file_path, filename)
    if suffix == ".csv":
        return create_csv_index(file_path, filename)
    if suffix in [".xlsx", ".xls"]:
        return create_excel_llama_index(file_path, filename)

    raise ValueError("Unsupported file type. Only .pdf, .docx, .pptx, .csv, .xlsx, .xls are supported.")


# -------------------------------
# STATE
# -------------------------------
class UploadRAGState(TypedDict):
    question: str
    file_index: str
    context: str
    approved: bool
    final_answer: str
    cached: bool
    context_precision: float
    sources: List[dict]
    chat_history:List[Dict[str,str]]#why lis tof dict
    applied_filters: Dict


def make_thread_id(file_index: str, session_id: str) -> str:
    raw = f"{session_id}|{file_index}|upload_chat"
    return hashlib.sha256(raw.encode()).hexdigest()

def format_history(history: List[Dict[str, str]], max_turns: int = 8) -> str:
    if not history:
        return ""
    tail = history[-2 * max_turns:]  # user+assistant pairs
    lines = []
    for m in tail:
        role = m.get("role", "")
        content = (m.get("content") or "").strip()
        if not content:
            continue
        if role == "user":
            lines.append(f"User: {content}")
        elif role == "assistant":
            lines.append(f"Assistant: {content}")
        else:
            lines.append(f"{role}: {content}")
    return "\n".join(lines).strip()



# -------------------------------
# RETRIEVAL (SYNC)
# -------------------------------
# REPLACE retrieve_sync with this
def retrieve_sync(state: UploadRAGState):
    start = datetime.now()
    file_id = state["file_index"]

    # Excel: LlamaIndex handles its own retrieval
    if file_id in _LLAMA_CACHE:
        query_engine = _LLAMA_CACHE[file_id]
        resp = query_engine.query(state["question"])
        context = str(resp)
        sources = [
            {
                "paper_id": file_id,
                "chunk_id": None,
                "source_type": "excel",
                "section_title": None,
                "preview": (context[:200] + "...") if len(context) > 200 else context,
            }
        ]
        logger.info("⏱ retrieve_sync (excel): %.4fs", (datetime.now() - start).total_seconds())
        return {"context": context, "sources": sources, "approved": True, "applied_filters": {}}

    # Derive source_type from filename — always known, no LLM needed
    suffix = Path(file_id).suffix.lower()
    source_type_map = {
        ".pdf": "pdf", ".docx": "docx", ".pptx": "pptx",
        ".csv": "csv", ".xlsx": "excel", ".xls": "excel",
    }
    source_type = source_type_map.get(suffix)

    filter_dict = {"file_id": file_id}
    if source_type:
        filter_dict["source_type"] = source_type

    qvec = EMBEDDINGS.embed_query(state["question"])
    res = query_embeddings(
        embedding=qvec,
        top_k=3,
        filter=filter_dict,
        min_score=0.3,  # skip irrelevant chunks
    )

    matches = res.get("matches", []) or []
    texts: List[str] = []
    sources: List[dict] = []

    for m in matches:
        md = (m.get("metadata") or {})
        txt = md.get("text") or ""
        md_small = dict(md)
        md_small.pop("text", None)
        if txt:
            texts.append(txt)
        sources.append(
            {
                **md_small,
                "preview": (txt[:200] + "...") if len(txt) > 200 else txt,
                "score": m.get("score"),
            }
        )

    context = "\n\n".join(texts)
    approved = bool(texts)  # False if no chunks passed min_score
    logger.info("⏱ retrieve_sync (qdrant): %.4fs", (datetime.now() - start).total_seconds())
    return {"context": context, "sources": sources, "approved": approved, "applied_filters": filter_dict}



# -------------------------------
# PROMPT
# -------------------------------
PROMPT = PromptTemplate(
    input_variables=["history", "context", "question"],
    template="""
You are a helpful assistant.
Use the conversation history for continuity, and answer using ONLY the context below.
If the answer is not present in the context, say exactly: "answer not found".

Conversation History:
{history}

Context:
{context}

Question:
{question}
""",
)


# -------------------------------
# RAGAS / CONTEXT PRECISION
# -------------------------------
_evaluator_llm = None
if OPENAI_API_KEY:
    try:
        client = OpenAI(api_key=OPENAI_API_KEY)
        _evaluator_llm = llm_factory("gpt-4o-mini", client=client)
    except Exception as e:
        logger.warning(f"Failed to init RAGAS evaluator LLM: {e}")
        _evaluator_llm = None
else:
    logger.warning("OPENAI_API_KEY not found; RAGAS context precision will be disabled.")
    _evaluator_llm = None


def compute_context_precision_sync(question: str, final_answer: str, context: str) -> float:
    if not _evaluator_llm:
        return 0.0
    metric = LLMContextPrecisionWithoutReference(llm=_evaluator_llm)
    sample = SingleTurnSample(
        user_input=question,
        response=final_answer,
        retrieved_contexts=[context],
    )
    score = metric.single_turn_score(sample)
    return float(score)


# -------------------------------
# GENERATE (SYNC STREAMING via LangGraph ONLY)
# -------------------------------
def generate_sync(state: UploadRAGState, config):
    if not state.get("approved"):
        yield {"final_answer": "answer not found", "context_precision": 0.0, "sources": []}
        return

    start = datetime.now()

    MAX_CONTEXT_CHARS = 8000
    trimmed_context = (state.get("context") or "")[:MAX_CONTEXT_CHARS]

    history_text = format_history(state.get("chat_history", []), max_turns=8)

    prompt = PROMPT.format(
        history=history_text,
        context=trimmed_context,
        question=state["question"],
    )

    parts: List[str] = []
    for chunk in LLM.stream(prompt):
        text = getattr(chunk, "content", None)
        if text:
            parts.append(text)
            yield {"final_answer": "".join(parts)}

    final_answer = "".join(parts).strip()

    score = compute_context_precision_sync(
        state["question"],
        final_answer,
        state.get("context") or "",
    )

    # ✅ IMPORTANT: write the new turn into state so MemorySaver persists it
    history = list(state.get("chat_history", []))
    history.append({"role": "user", "content": state["question"]})
    history.append({"role": "assistant", "content": final_answer})

    logger.info("⏱ generate_sync (LLM stream): %.4fs", (datetime.now() - start).total_seconds())
    logger.info("context precision computed (upload): %s", score)

    yield {
        "final_answer": final_answer,
        "context_precision": score,
        "sources": state.get("sources", []),
        "chat_history": history,  # ✅ persisted by checkpointer
    }


# -------------------------------
# GRAPH (LangGraph) - SYNC
# -------------------------------
graph = StateGraph(UploadRAGState)
graph.add_node("retrieve", retrieve_sync)
graph.add_node("generate", generate_sync)
graph.set_entry_point("retrieve")
graph.add_edge("retrieve", "generate")
graph.add_edge("generate", END)
app = graph.compile(checkpointer=MemorySaver())


def _lg_config(thread_id: str):
    return {"configurable": {"thread_id": thread_id}}


# -------------------------------
# PUBLIC API: NON-STREAM (sync final)
# -------------------------------
def answer_from_uploaded_file(question: str, file_index: str, session_id: str) -> dict:
    cache_key = make_cache_key(question, file_index)
    cached = cache_get(cache_key)
    if cached:
        answer, sources, cp = cached
        return {"final_answer": answer, "cached": True, "sources": sources, "context_precision": float(cp)}

    state: UploadRAGState = {
        "question": question,
        "file_index": file_index,
        "context": "",
        "approved": False,
        "final_answer": "",
        "cached": False,
        "context_precision": 0.0,
        "sources": [],
        "chat_history": [],  # ✅ init
        "applied_filters": {}
    }

    thread_id = make_thread_id(file_index, session_id)
    final_state = app.invoke(state, config=_lg_config(thread_id=thread_id))
    if not isinstance(final_state, dict):
        final_state = dict(final_state)

    cache_set(
        cache_key,
        final_state.get("final_answer", ""),
        final_state.get("sources", []),
        context_precision=float(final_state.get("context_precision", 0.0)),
    )

    final_state["cached"] = False
    return final_state


# -------------------------------
# PUBLIC API: STREAM (LangGraph only) - SYNC
# -------------------------------
def stream_answer_from_uploaded_file(question: str, file_index: str, session_id: str):
    cache_key = make_cache_key(question, file_index)
    cached = cache_get(cache_key)
    if cached:
        answer, sources, cp = cached
        yield {"event": "cached", "final_answer": answer, "sources": sources, "context_precision": float(cp), "cached": True}
        return

    state: UploadRAGState = {
        "question": question,
        "file_index": file_index,
        "context": "",
        "approved": False,
        "final_answer": "",
        "cached": False,
        "context_precision": 0.0,
        "sources": [],
        "chat_history": [],  # ✅ init (MemorySaver will replace with stored one automatically)
        "applied_filters": {}
    }

    thread_id = make_thread_id(file_index, session_id)

    last_answer = ""
    last_sources = []
    last_cp = 0.0

    for upd in app.stream(state, stream_mode="updates", config=_lg_config(thread_id=thread_id)):
        if not upd:
            continue

        if "retrieve" in upd:
            r = upd["retrieve"] or {}
            last_sources = r.get("sources", []) or last_sources
            last_filters = r.get("applied_filters", {})
            yield {"event": "retrieved", "sources": last_sources, "applied_filters": last_filters}

        if "generate" in upd:
            g = upd["generate"] or {}
            if "final_answer" in g and g["final_answer"] is not None:
                last_answer = g["final_answer"]
                yield {"event": "answer", "final_answer": last_answer}
            if "context_precision" in g and g["context_precision"] is not None:
                last_cp = float(g["context_precision"])
            if "sources" in g and g["sources"] is not None:
                last_sources = g["sources"]

    if last_answer:
        cache_set(cache_key, last_answer, last_sources, context_precision=float(last_cp))

    yield {"event": "done", "final_answer": last_answer, "sources": last_sources, "context_precision": float(last_cp), "cached": False}


def answer_from_uploaded_pdf(question: str, pdf_index: str) -> dict:
    return answer_from_uploaded_file(question=question, file_index=pdf_index)
